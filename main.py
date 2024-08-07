import os
import wx
from pptx import Presentation


class MergePPTFrame(wx.Frame):
    def __init__(self, parent, title):
        super(MergePPTFrame, self).__init__(parent, title=title, size=(400, 200))

        panel = wx.Panel(self)
        vbox = wx.BoxSizer(wx.VERTICAL)

        self.folder_path = wx.StaticText(panel, label="选择文件夹:")
        vbox.Add(self.folder_path, flag=wx.ALL, border=10)

        self.merge_button = wx.Button(panel, label="合并PPT")
        self.merge_button.Bind(wx.EVT_BUTTON, self.on_merge)
        vbox.Add(self.merge_button, flag=wx.ALIGN_CENTER | wx.ALL, border=10)

        panel.SetSizer(vbox)
        self.Show(True)

    def on_merge(self, event):
        dialog = wx.DirDialog(self, "选择文件夹", style=wx.DD_DEFAULT_STYLE | wx.DD_DIR_MUST_EXIST)
        if dialog.ShowModal() == wx.ID_OK:
            folder_path = dialog.GetPath()
            self.folder_path.SetLabelText("选择文件夹: {}".format(folder_path))
            self.merge_ppt_files(folder_path)
        dialog.Destroy()

    def merge_ppt_files(self, folder_path):
        output_ppt = Presentation()

        for root, dirs, files in os.walk(folder_path):
            for file in files:
                if file.endswith(".ppt") or file.endswith(".pptx"):
                    ppt_path = os.path.join(root, file)
                    presentation = Presentation(ppt_path)
                    for slide in presentation.slides._sldIdLst:
                        output_ppt.slides._sldIdLst.append(slide)

        output_ppt.save("merged_ppt.pptx")
        wx.MessageBox("PPT文件合并完成！", "提示", wx.OK | wx.ICON_INFORMATION)


app = wx.App()
MergePPTFrame(None, title='PPT合并工具')
app.MainLoop()